from pathlib import Path
from typing import List
from sqlalchemy.orm import Session
import logging
from backend.celery_app import celery_app
from backend.database import SessionLocal
from backend.models import Document, DocumentChunk, Workspace
from backend.services.embeddings import get_embeddings_model
from backend.services.ingestion import (
    promote_structural_markers,
    chunk_markdown,
    batch_iter,
    EMBED_BATCH_SIZE,
)

# Multimodal loaders
import pymupdf4llm
from docx import Document as DocxDocument
from pptx import Presentation
import openai

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


@celery_app.task(name="process_document_task")
def process_document_task(document_id: int):
    """
    Asynchronous task to process a single document based on its file extension.
    """
    db = SessionLocal()
    db_doc = db.query(Document).filter(Document.id == document_id).first()

    if not db_doc:
        logger.error(f"Document {document_id} not found.")
        return

    try:
        db_doc.status = "processing"
        db.commit()

        file_path = Path(db_doc.file_path)
        ext = file_path.suffix.lower()

        content_pages = []  # List of {"text": str, "metadata": dict}

        if ext == ".pdf":
            content_pages = pymupdf4llm.to_markdown(str(file_path), page_chunks=True)
            db_doc.file_type = "pdf"
        elif ext in [".docx", ".doc"]:
            content_pages = _process_docx(file_path)
            db_doc.file_type = "docx"
        elif ext in [".pptx", ".ppt"]:
            content_pages = _process_pptx(file_path)
            db_doc.file_type = "pptx"
        elif ext in [".jpg", ".jpeg", ".png", ".webp"]:
            content_pages = _process_image(file_path, db, db_doc.workspace_id)
            db_doc.file_type = "image"
        else:
            raise ValueError(f"Unsupported file extension: {ext}")

        # Finalize and store chunks
        total_chunks = store_processed_content(db, db_doc, content_pages)

        db_doc.status = "completed"
        db.commit()
        logger.info(
            f"Successfully processed document {document_id} ({total_chunks} chunks)"
        )

    except Exception as e:
        logger.exception(f"Error processing document {document_id}")
        db_doc.status = "failed"
        db_doc.error_message = str(e)
        db.commit()
    finally:
        db.close()


def _process_docx(path: Path) -> List[dict]:
    doc = DocxDocument(str(path))
    full_text = []
    for para in doc.paragraphs:
        full_text.append(para.text)

    # Simple page-like chunking for Word
    return [{"text": "\n".join(full_text), "metadata": {"page": 1}}]


def _process_pptx(path: Path) -> List[dict]:
    prs = Presentation(str(path))
    content_pages = []
    for i, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        content_pages.append(
            {"text": "\n".join(slide_text), "metadata": {"page": i + 1}}
        )
    return content_pages


def _process_image(path: Path, db: Session, workspace_id: int) -> List[dict]:
    """
    Uses Vision LLM to describe the image.
    Supports both OpenAI Vision and Ollama Vision (llava/bakllava).
    """
    import base64
    import requests
    from backend.services.settings import get_app_settings

    settings_db = get_app_settings(db)

    # Check if vision processing is enabled
    if not settings_db.enable_vision_processing:
        raise ValueError(
            "Image processing is disabled in settings. "
            "Enable vision processing in Settings or upload a PDF/Word/PPT instead."
        )

    prompt_text = "Describe this image in extreme detail for an educational RAG system. Extract all text, explain diagrams, and summarize key concepts shown."

    with open(path, "rb") as image_file:
        encoded_string = base64.b64encode(image_file.read()).decode("utf-8")

    vision_provider = settings_db.vision_provider or "openai"

    # Ollama Vision Support
    if vision_provider == "ollama":
        try:
            ollama_url = settings_db.ollama_base_url or "http://localhost:11434"
            vision_model = settings_db.ollama_vision_model or "llava"

            response = requests.post(
                f"{ollama_url}/api/generate",
                json={
                    "model": vision_model,
                    "prompt": prompt_text,
                    "images": [encoded_string],
                    "stream": False,
                },
                timeout=120,  # Vision models can be slow
            )
            response.raise_for_status()
            description = response.json().get("response", "")

            if not description:
                raise ValueError(
                    f"Ollama vision model '{vision_model}' returned empty response"
                )

        except Exception as e:
            logger.error(f"Ollama vision processing failed: {e}")
            raise ValueError(
                f"Ollama vision processing failed: {str(e)}. "
                f"Make sure '{vision_model}' model is installed (run: ollama pull {vision_model})"
            )

    # OpenAI Vision Support
    else:
        if not settings_db.openai_api_key:
            raise ValueError(
                "Image processing requires OpenAI Vision, but no OpenAI API key is configured. "
                "Add your key in Settings, switch to Ollama vision, or upload a PDF/Word/PPT instead."
            )

        workspace = db.get(Workspace, workspace_id)
        vision_model = settings_db.openai_model or "gpt-4o"
        if workspace and workspace.llm_provider == "openai" and workspace.llm_model:
            vision_model = workspace.llm_model

        client = openai.OpenAI(api_key=settings_db.openai_api_key)
        vision_response = client.chat.completions.create(
            model=vision_model,
            messages=[
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": prompt_text},
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/jpeg;base64,{encoded_string}"
                            },
                        },
                    ],
                }
            ],
            max_tokens=1500,
        )
        description = vision_response.choices[0].message.content

    return [
        {
            "text": description or "",
            "metadata": {"page": 1, "type": "image_description"},
        }
    ]


def store_processed_content(db, db_doc, pages: List[dict]) -> int:
    model, dim, provider, model_name = get_embeddings_model(db, db_doc.workspace_id)

    # Record the model used for this document
    db_doc.embedding_provider = provider
    db_doc.embedding_model = model_name
    db.commit()
    all_rows: List[DocumentChunk] = []
    chunk_index = 0

    for page_data in pages:
        page_text = page_data.get("text", "")
        page_num = page_data.get("metadata", {}).get("page", 1)

        refined = promote_structural_markers(page_text)
        chunks = chunk_markdown(refined)

        for _, batch in batch_iter(chunks, EMBED_BATCH_SIZE):
            texts = [c.page_content for c in batch]
            vectors = model.embed_documents(texts)

            for i, (chunk, vector) in enumerate(zip(batch, vectors)):
                meta = chunk.metadata.copy()
                meta["page"] = page_num
                meta["source"] = db_doc.title

                # Context prefix logic from ingestion.py
                headers = [
                    str(meta.get(f"Header {j}"))
                    for j in range(1, 7)
                    if meta.get(f"Header {j}")
                ]
                prefix = f"Context: {' > '.join(headers) if headers else db_doc.title} (Page {page_num})"
                enriched_content = f"{prefix}\n\n{chunk.page_content}"

                chunk_args = {
                    "document_id": db_doc.id,
                    "workspace_id": db_doc.workspace_id,
                    "content": enriched_content,
                    "chunk_index": chunk_index,
                    "chunk_metadata": meta,
                }

                # Assign to correct embedding column
                if dim == 1536:
                    chunk_args["embedding_1536"] = vector
                elif dim == 1024:
                    chunk_args["embedding_1024"] = vector
                elif dim == 768:
                    chunk_args["embedding_768"] = vector
                elif dim == 384:
                    chunk_args["embedding_384"] = vector
                else:
                    chunk_args["embedding_768"] = vector

                all_rows.append(DocumentChunk(**chunk_args))
                chunk_index += 1

    db.add_all(all_rows)
    db.commit()
    return chunk_index
